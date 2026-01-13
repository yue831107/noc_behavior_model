#!/usr/bin/env python3
"""
Generate NoC Performance Analysis Presentation.

Creates a PowerPoint presentation covering:
1. Performance validation architecture
2. Hardware tunable parameters
3. Parameter optimization methodology

Usage:
    py -3 tools/gen_presentation.py
    py -3 tools/gen_presentation.py -o output/my_presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pathlib import Path
import argparse


def add_title_slide(prs, title, subtitle=""):
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(2.5)
    width = Inches(9)
    height = Inches(1.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    if subtitle:
        top = Inches(4)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_slide(prs, section_title, section_number=""):
    """Add a section divider slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Section number
    if section_number:
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(0.8)

        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Part {section_number}"
        p.font.size = Pt(20)
        # Gray color for section number
        from pptx.dml.color import RGBColor
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.alignment = PP_ALIGN.CENTER

    # Section title
    left = Inches(0.5)
    top = Inches(2.8)
    width = Inches(9)
    height = Inches(1.2)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = section_title
    p.font.size = Pt(36)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_content_slide(prs, title, bullet_points, notes=""):
    """Add a content slide with bullet points."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Bullet points
    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(9)
    height = Inches(5.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Handle indentation
        if isinstance(point, tuple):
            text, level = point
            p.text = text
            p.level = level
        else:
            p.text = f"• {point}"
            p.level = 0

        p.font.size = Pt(18)
        p.space_after = Pt(8)

    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes

    return slide


def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Table
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)

    left = Inches(0.5)
    top = Inches(1.3)
    width = Inches(9)
    height = Inches(0.4 * num_rows)

    table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

    # Set column widths
    col_width = Inches(9 / num_cols)
    for i in range(num_cols):
        table.columns[i].width = col_width

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_data)
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide


def add_code_slide(prs, title, code_text):
    """Add a slide with code block."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.7)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    # Code block
    left = Inches(0.5)
    top = Inches(1.1)
    width = Inches(9)
    height = Inches(5.5)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(11)
    p.font.name = "Consolas"

    return slide


def create_presentation(output_path: Path):
    """Create the full presentation."""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # =========================================================================
    # Title Slide
    # =========================================================================
    add_title_slide(
        prs,
        "NoC Behavior Model",
        "Performance Validation & Parameter Optimization"
    )

    # =========================================================================
    # Part 1: Performance Validation Architecture
    # =========================================================================
    add_section_slide(prs, "Performance Validation Architecture", "1")

    # Overview
    add_content_slide(prs, "Validation Architecture Overview", [
        "Monitor-Based Architecture: External observation, no core modification",
        "Two-Layer Validation:",
        ("  • Theory Validation: Throughput, Latency bounds", 1),
        ("  • Consistency Validation: Little's Law, Flit Conservation", 1),
        "Key Principle: Validator depends on Core (one-way dependency)",
        "MetricsCollector: Unified interface for all metrics",
    ])

    # Validation Flow
    add_content_slide(prs, "Validation Flow", [
        "1. NoC Core executes simulation",
        "2. MetricsProvider extracts metrics (zero intrusion)",
        "3. MetricsCollector captures time-series snapshots",
        "4. Validators check against theoretical bounds",
        "5. Report generation with PASS/FAIL status",
    ])

    # Theory Validation
    add_table_slide(prs, "Theory Validation Formulas",
        ["Metric", "Formula", "5x4 Mesh Value"],
        [
            ["Throughput Max", "T_max = N_edge × W_flit", "4 × 8 = 32 B/cycle"],
            ["Latency Min", "L_min = hops × P_depth + 2", "7 cycles (fast mode)"],
            ["Buffer Util", "U = avg_occupancy / capacity", "0 ≤ U ≤ 1"],
        ]
    )

    # Consistency Validation
    add_table_slide(prs, "Consistency Validation",
        ["Check", "Formula", "Tolerance"],
        [
            ["Little's Law", "L = λ × W", "10%"],
            ["Flit Conservation", "Sent = Received", "0%"],
            ["Bandwidth Conservation", "Σ(inject) = Σ(eject)", "10%"],
        ]
    )

    # Latency Tracking
    add_content_slide(prs, "Monitor-Based Latency Tracking", [
        "Problem: Core should not embed monitoring logic",
        "Solution: External injection/ejection tracking",
        "",
        "API:",
        ("  • record_injection(key, cycle) - When request starts", 1),
        ("  • record_ejection(key, cycle) - When response arrives", 1),
        ("  • latency = ejection_cycle - injection_cycle", 1),
        "",
        "Supports both Host-to-NoC (AXI ID) and NoC-to-NoC (Node ID)",
    ])

    # Host vs NoC Validation Differences
    add_table_slide(prs, "Host-to-NoC vs NoC-to-NoC Validation",
        ["Validation", "Host-to-NoC", "NoC-to-NoC", "Reason"],
        [
            ["Throughput Bound", "✓ Validate", "✗ Skip", "No edge bottleneck"],
            ["Little's Law", "✓ Validate", "✗ Skip", "Burst traffic"],
            ["Buffer Util", "✓ Validate", "✓ Validate", "Both applicable"],
            ["Flit Conservation", "✓ Validate", "✓ Validate", "Both applicable"],
        ]
    )

    # =========================================================================
    # Part 2: Hardware Tunable Parameters
    # =========================================================================
    add_section_slide(prs, "Hardware Tunable Parameters", "2")

    # Parameter Overview
    add_content_slide(prs, "Parameter Categories", [
        "Mesh Topology Parameters:",
        ("  • cols × rows: Network size (fixed at 5×4)", 1),
        ("  • edge_column: Edge router column (fixed at 0)", 1),
        "",
        "Router Parameters:",
        ("  • buffer_depth: Input buffer size per port [2-32]", 1),
        ("  • pipeline_mode: fast/standard/hardware", 1),
        "",
        "NI Parameters:",
        ("  • max_outstanding: Concurrent transactions [4-64]", 1),
        ("  • flit_payload_size: Bytes per flit [16-128]", 1),
    ])

    # Buffer Depth
    add_table_slide(prs, "Buffer Depth Impact",
        ["buffer_depth", "Congestion Tolerance", "Area Cost", "Use Case"],
        [
            ["2-4", "Low", "Low", "Low load, area sensitive"],
            ["4-8", "Medium", "Medium", "General purpose"],
            ["8-16", "High", "High", "High load, burst traffic"],
            ["16-32", "Very High", "Very High", "Extreme congestion"],
        ]
    )

    # Max Outstanding
    add_table_slide(prs, "Max Outstanding Impact",
        ["max_outstanding", "Throughput Potential", "Memory Cost", "Use Case"],
        [
            ["4-8", "Low", "Low", "Simple control flow"],
            ["16", "Medium", "Medium", "General purpose"],
            ["32-64", "High", "High", "High throughput demand"],
        ]
    )

    # Parameter Interaction Matrix
    add_table_slide(prs, "Parameter Impact Matrix",
        ["Parameter ↑", "Throughput", "Latency", "Congestion Tolerance"],
        [
            ["buffer_depth", "○ No effect", "✗ Slight increase", "✓ Major improvement"],
            ["max_outstanding", "✓ Significant", "○ No effect", "○ No effect"],
            ["mesh_size", "✓ Capacity up", "✗ Distance up", "○ No effect"],
        ]
    )

    # Critical Discovery
    add_content_slide(prs, "Critical Discovery: Buffer-Outstanding Mismatch", [
        "Problem Found: Deadlock when buffer_depth < max_outstanding",
        "",
        "Root Cause:",
        ("  • max_outstanding = 16 (fixed in current design)", 1),
        ("  • selector_ingress_buffer = buffer_depth", 1),
        ("  • When buffer_depth=4: only 4 flits can be buffered", 1),
        ("  • Remaining 12 transactions block → Deadlock!", 1),
        "",
        "Solution: Ensure buffer_depth ≥ max_outstanding",
        "",
        "Evidence: buffer_depth=16 achieves 28.85 B/cycle",
        "          buffer_depth=4 achieves only 0.08 B/cycle",
    ])

    # =========================================================================
    # Part 3: Parameter Optimization
    # =========================================================================
    add_section_slide(prs, "Parameter Optimization Methodology", "3")

    # Optimization Flow
    add_content_slide(prs, "Optimization Flow", [
        "1. Define Performance Targets",
        ("  • Throughput: minimum acceptable B/cycle", 1),
        ("  • Latency: maximum acceptable cycles", 1),
        ("  • Buffer Utilization: maximum acceptable ratio", 1),
        "",
        "2. Define Parameter Search Space",
        ("  • buffer_depth: [2, 4, 8, 16, 32]", 1),
        ("  • max_outstanding: [4, 8, 16, 32]", 1),
        "",
        "3. Execute Parameter Sweep",
        "",
        "4. Evaluate & Rank Configurations",
    ])

    # Performance Targets
    add_table_slide(prs, "Performance Targets (5×4 Mesh)",
        ["Mode", "Throughput", "Latency (avg)", "Latency (max)", "Buffer Util"],
        [
            ["Host-to-NoC", "≥ 20 B/cycle", "≤ 150 cycles", "≤ 1000 cycles", "≤ 70%"],
            ["NoC-to-NoC", "≥ 80 B/cycle", "≤ 50 cycles", "≤ 200 cycles", "≤ 70%"],
        ]
    )

    # Scoring Function
    add_content_slide(prs, "Configuration Scoring Function", [
        "Score = Throughput_Score + Latency_Score + Buffer_Score",
        "",
        "Throughput Score (40% weight):",
        ("  • Meet target: 40 × (1 + excess_ratio)", 1),
        ("  • Miss target: 40 × (1 - deficit_ratio)", 1),
        "",
        "Latency Score (40% weight):",
        ("  • Meet target: 40 × (1 + margin_ratio)", 1),
        ("  • Miss target: 40 × (1 - excess_ratio)", 1),
        "",
        "Buffer Score (20% weight):",
        ("  • Lower utilization = higher score", 1),
    ])

    # Sweep Results - Host to NoC
    add_table_slide(prs, "Sweep Results: Host-to-NoC",
        ["buffer_depth", "Throughput", "Latency (avg)", "Status"],
        [
            ["2", "0.08 B/cycle", "19 cycles", "✗ FAIL (deadlock)"],
            ["4", "0.08 B/cycle", "29 cycles", "✗ FAIL (deadlock)"],
            ["8", "0.08 B/cycle", "45 cycles", "✗ FAIL (deadlock)"],
            ["16", "28.85 B/cycle", "77 cycles", "✓ PASS"],
        ]
    )

    # Sweep Results - NoC to NoC
    add_table_slide(prs, "Sweep Results: NoC-to-NoC",
        ["buffer_depth", "Throughput", "Latency (avg)", "Status"],
        [
            ["2", "163.84 B/cycle", "16.8 cycles", "✓ PASS"],
            ["4", "163.84 B/cycle", "16.8 cycles", "✓ PASS"],
            ["8", "163.84 B/cycle", "16.8 cycles", "✓ PASS"],
            ["16", "163.84 B/cycle", "16.8 cycles", "✓ PASS"],
        ]
    )

    # Optimal Configuration
    add_table_slide(prs, "Optimal Configuration Recommendation",
        ["Mode", "buffer_depth", "max_outstanding", "Rationale"],
        [
            ["Host-to-NoC", "16", "16", "Minimum to avoid deadlock"],
            ["NoC-to-NoC", "2", "4", "Minimum that meets target"],
        ]
    )

    # Tool Usage
    add_code_slide(prs, "Parameter Sweep Tool Usage", """
# Run parameter sweep for both modes
py -3 tools/param_sweep.py --mode both

# Custom parameter ranges
py -3 tools/param_sweep.py \\
    --buffer-depths 4 8 16 32 \\
    --max-outstandings 8 16 32 64

# Output files:
#   output/param_sweep/sweep_host_to_noc.json
#   output/param_sweep/sweep_noc_to_noc.json

# Batch performance test (500+ configs)
py -3 tools/run_batch_perf_test.py --mode both --count 500
""")

    # Key Takeaways
    add_content_slide(prs, "Key Takeaways", [
        "1. Monitor-based validation: Don't modify core code",
        "",
        "2. Buffer-Outstanding balance is critical:",
        ("  • buffer_depth must be ≥ max_outstanding", 1),
        ("  • Mismatch causes deadlock in Host-to-NoC", 1),
        "",
        "3. NoC-to-NoC is more robust:",
        ("  • No edge router bottleneck", 1),
        ("  • All configurations meet target", 1),
        "",
        "4. Use automated sweep tools for optimization",
    ])

    # =========================================================================
    # Save
    # =========================================================================
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


def main():
    parser = argparse.ArgumentParser(
        description='Generate NoC Performance Analysis Presentation'
    )
    parser.add_argument(
        '-o', '--output',
        type=Path,
        default=Path('output/NoC_Performance_Analysis.pptx'),
        help='Output path (default: output/NoC_Performance_Analysis.pptx)'
    )

    args = parser.parse_args()
    create_presentation(args.output)


if __name__ == '__main__':
    main()
